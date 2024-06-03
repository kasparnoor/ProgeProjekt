import pptx
from pptx.util import Inches

# Load an existing presentation (or create a new one)
prs = pptx.Presentation("my_presentation.pptx")  # Or pptx.Presentation() for new

# Select a slide (you can use slide index or layout)
slide = prs.slides[0]  # Get the first slide

# Add an image
left = Inches(1)
top = Inches(2)
pic = slide.shapes.add_picture("my_image.jpg", left, top)

# Optional: Adjust image size
pic.width = Inches(4)
pic.height = Inches(3)

# Save the presentation
prs.save("my_presentation_updated.pptx")