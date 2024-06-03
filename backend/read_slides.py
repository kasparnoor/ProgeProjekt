import pptx

def extract_text_from_slides(filepath):
    prs = pptx.Presentation(filepath)
    text_list = []

    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text += run.text
        text_list.append(slide_text)

    return text_list


if __name__ == "__main__":
    presentation_file = "my_presentation.pptx"
    all_slide_text = extract_text_from_slides(presentation_file)

    print(all_slide_text)  # This will print a list of text from each slide
